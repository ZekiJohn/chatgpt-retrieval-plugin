import os
from io import BufferedReader
from typing import Optional
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader, PdfReader
import docx2txt
import csv
import pptx
from models.models import Document, DocumentMetadata, Source
from io import BytesIO


async def get_document_from_file(file: bytes, mimetype: str) -> Document:
    extracted_text = await extract_text_from_form_file(file, mimetype)
    metadata = DocumentMetadata(
        source=Source.file,
    )
    doc = Document(text=extracted_text, metadata=metadata)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    # Open the file in binary mode
    file = open(filepath, "rb")
    extracted_text = extract_text_from_file(file, mimetype)

    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        print(file)
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        file.close()
        raise ValueError("Unsupported file type: {}".format(mimetype))

    file.close()
    return extracted_text


# Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: bytes, mimetype: str):
    """Return the text content of a file."""
    # get the file body from the upload file object

    # mimetype = file.content_type

    # print(f"mimetype: {mimetype}")
    # print(f"file.file: {file.file}")
    # print("file: ", file)

    file_stream = file
    # file_stream = await file.read()

    temp_file_path = "/tmp/temp_file"

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)

    try:
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)
    except Exception as e:
        print(f"Error: {e}")
        os.remove(temp_file_path)
        raise e

    # remove file from temp location
    os.remove(temp_file_path)

    return extracted_text



async def count_characters_in_pdf(file_stream: bytes, mimetype: str) -> int:
    # if mimetype == "application/pdf":
    # elif mimetype == "text/plain" or mimetype == "text/markdown":
    # elif (
    #         mimetype
    #         == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    #     ):
    # elif mimetype == "text/csv":
    # elif (
    #         mimetype
    #         == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    #     ):
    # else:
    #     # Unsupported file type

    if not file_stream:
        raise ValueError("The provided file is empty")
    # Create a PDF file reader
    with BytesIO(file_stream) as pdf_file:
        pdf_reader = PdfReader(pdf_file)
        num_pages = len(pdf_reader.pages)
        
        total_chars = 0

        # Iterate through all the pages and count the characters
        for i in range(num_pages):
            page = pdf_reader.pages[i]
            text = page.extract_text()
            total_chars += len(text)
    
    return total_chars


# def count_characters(file_path):
#     with open(file_path, 'rb') as f:
#         pdf_reader = PdfReader(f)
#         num_pages = pdf_reader.getNumPages()
#         char_count = 0
#         for page_num in range(num_pages):
#             page_obj = pdf_reader.pages[page_num]
#             page_text = page_obj.extract_text()
#             char_count += len(page_text)
#     return char_count

# def count_characters_in_pdf(pdf_content: bytes) -> int:
#     with BytesIO(pdf_content) as pdf_file:
#         text = extract_text(pdf_file)
#     return len(text)
